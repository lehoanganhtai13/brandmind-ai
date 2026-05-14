"""Tests for artifact audit provenance selection."""

from __future__ import annotations

import json
from pathlib import Path

from evaluation.artifact_audit import (
    _check_presentation_semantic,
    audit,
    parse_metadata,
)


def test_parse_metadata_reads_api_session_id_fallback(tmp_path: Path) -> None:
    """A partial pilot directory still carries the API session id."""
    (tmp_path / "api_session_id.txt").write_text("api123\n", encoding="utf-8")

    api_session_id, scope, completed = parse_metadata(tmp_path)

    assert api_session_id == "api123"
    assert scope is None
    assert completed == []


def test_audit_uses_workspace_session_id_for_manifest(tmp_path: Path) -> None:
    """Manifest records are keyed by workspace session id, not API session id."""
    session_dir = tmp_path / "eval" / "pilot"
    session_dir.mkdir(parents=True)
    (session_dir / "api_session_id.txt").write_text("api123\n", encoding="utf-8")

    brandmind_home = tmp_path / "brandmind-home"
    workspace_dir = brandmind_home / "projects" / "workspace123" / "workspace"
    workspace_dir.mkdir(parents=True)
    (workspace_dir.parent / "project.json").write_text(
        json.dumps({"session_id": "workspace123", "brand_name": "Demo"}),
        encoding="utf-8",
    )
    (workspace_dir / "brand_brief.md").write_text("# Demo\n", encoding="utf-8")

    output_root = tmp_path / "brandmind-output"
    output_root.mkdir()
    doc_path = output_root / "documents" / "demo.docx"
    doc_path.parent.mkdir()
    doc_path.write_text("placeholder", encoding="utf-8")
    stale_path = output_root / "documents" / "stale.docx"
    stale_path.write_text("stale", encoding="utf-8")
    other_brand_path = output_root / "documents" / "other-brand.docx"
    other_brand_path.write_text("other", encoding="utf-8")
    manifest_rows = [
        {
            "session_id": "workspace123",
            "brand_name": "Demo",
            "category": "documents",
            "tool": "generate_document",
            "filename": doc_path.name,
            "path": str(doc_path),
        },
        {
            "session_id": "other-session",
            "brand_name": "Other",
            "category": "documents",
            "tool": "generate_document",
            "filename": stale_path.name,
            "path": str(stale_path),
        },
        {
            "session_id": "workspace123",
            "brand_name": "Other Brand",
            "category": "documents",
            "tool": "generate_document",
            "filename": other_brand_path.name,
            "path": str(other_brand_path),
        },
    ]
    (output_root / ".manifest.jsonl").write_text(
        "\n".join(json.dumps(row) for row in manifest_rows),
        encoding="utf-8",
    )

    report = audit(
        session_dir=session_dir,
        brandmind_home=brandmind_home,
        output_root=output_root,
    )

    assert report.session_id_api == "api123"
    assert report.workspace_session_id == "workspace123"
    assert report.artifacts_on_disk["strategy_document"] == [str(doc_path)]


def test_audit_accepts_current_session_brand_label_suffixes(tmp_path: Path) -> None:
    """Current-session artifacts survive harmless brand-label drift."""
    session_dir = tmp_path / "eval" / "pilot"
    session_dir.mkdir(parents=True)
    (session_dir / "api_session_id.txt").write_text("api123\n", encoding="utf-8")

    brandmind_home = tmp_path / "brandmind-home"
    workspace_dir = brandmind_home / "projects" / "workspace123" / "workspace"
    workspace_dir.mkdir(parents=True)
    (workspace_dir.parent / "project.json").write_text(
        json.dumps(
            {"session_id": "workspace123", "brand_name": "Chuyện Ba Bữa Signature"}
        ),
        encoding="utf-8",
    )
    (workspace_dir / "brand_brief.md").write_text("# Demo\n", encoding="utf-8")

    output_root = tmp_path / "brandmind-output"
    doc_path = output_root / "documents" / "brand_strategy.docx"
    sheet_path = output_root / "spreadsheets" / "kpi_dashboard.xlsx"
    other_path = output_root / "documents" / "other_brand.docx"
    for path in (doc_path, sheet_path, other_path):
        path.parent.mkdir(parents=True, exist_ok=True)
        path.write_text("placeholder", encoding="utf-8")

    manifest_rows = [
        {
            "session_id": "workspace123",
            "brand_name": "ChuyỆn Ba Bữa Signature",
            "category": "documents",
            "tool": "generate_document",
            "filename": doc_path.name,
            "path": str(doc_path),
        },
        {
            "session_id": "workspace123",
            "brand_name": "Chuyện Ba Bữa Signature - KPI",
            "category": "spreadsheets",
            "tool": "generate_spreadsheet",
            "filename": sheet_path.name,
            "path": str(sheet_path),
        },
        {
            "session_id": "workspace123",
            "brand_name": "Other Brand",
            "category": "documents",
            "tool": "generate_document",
            "filename": other_path.name,
            "path": str(other_path),
        },
    ]
    (output_root / ".manifest.jsonl").write_text(
        "\n".join(json.dumps(row, ensure_ascii=False) for row in manifest_rows),
        encoding="utf-8",
    )

    report = audit(
        session_dir=session_dir,
        brandmind_home=brandmind_home,
        output_root=output_root,
    )

    assert report.artifacts_on_disk["strategy_document"] == [str(doc_path)]
    assert report.artifacts_on_disk["spreadsheet"] == [str(sheet_path)]


def test_presentation_semantic_uses_first_text_when_title_placeholder_absent(
    tmp_path: Path,
) -> None:
    """Deck checker counts visible title text even without title placeholders."""
    from pptx import Presentation
    from pptx.util import Inches

    deck_path = tmp_path / "deck.pptx"
    presentation = Presentation()
    blank = presentation.slide_layouts[6]
    titles = [
        "Business Challenge",
        "Brand Positioning",
        "Brand Identity",
        "Communication Framework",
        "Implementation Roadmap",
        "KPIs & Measurement",
        "Executive Summary",
        "Appendix",
    ]
    for title in titles:
        slide = presentation.slides.add_slide(blank)
        box = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(8), Inches(1))
        box.text_frame.text = title
    presentation.save(deck_path)

    check = _check_presentation_semantic(deck_path)

    assert check.passed
    assert check.details["slide_count"] == 8
    assert check.details["title_topics_matched"] == 6


def test_audit_does_not_count_old_artifacts_without_provenance(
    tmp_path: Path,
) -> None:
    """Current-session audit must not pass from old same-brand files."""
    session_dir = tmp_path / "eval" / "pilot"
    session_dir.mkdir(parents=True)
    (session_dir / "metadata.json").write_text(
        json.dumps(
            {
                "session_id": "api123",
                "date": "20260512",
                "session_metadata": {
                    "scope": "repositioning",
                    "completed_phases": ["phase_0"],
                },
            }
        ),
        encoding="utf-8",
    )

    brandmind_home = tmp_path / "brandmind-home"
    workspace_dir = brandmind_home / "projects" / "workspace123" / "workspace"
    workspace_dir.mkdir(parents=True)
    (workspace_dir.parent / "project.json").write_text(
        json.dumps({"session_id": "workspace123", "brand_name": "Demo"}),
        encoding="utf-8",
    )
    (workspace_dir / "brand_brief.md").write_text(
        "# Demo\n\n## Phase 0: Business Problem Diagnosis\n",
        encoding="utf-8",
    )

    output_root = tmp_path / "brandmind-output"
    old_doc = output_root / "documents" / "demo" / "old_brand_strategy.docx"
    old_doc.parent.mkdir(parents=True)
    old_doc.write_text("old artifact", encoding="utf-8")
    (output_root / ".manifest.jsonl").write_text(
        json.dumps(
            {
                "session_id": "old-session",
                "brand_name": "Demo",
                "category": "documents",
                "tool": "generate_document",
                "filename": old_doc.name,
                "path": str(old_doc),
            }
        ),
        encoding="utf-8",
    )

    report = audit(
        session_dir=session_dir,
        brandmind_home=brandmind_home,
        output_root=output_root,
    )

    assert report.artifacts_on_disk["strategy_document"] == []
    assert not report.tier1_health.strategy_doc_produced
