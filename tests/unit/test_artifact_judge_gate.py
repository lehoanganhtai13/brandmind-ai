"""Tests for the strict artifact content acceptance gate."""

from __future__ import annotations

from evaluation.judge.artifact_judge import (
    _LEVEL_ACCEPTABLE,
    _LEVEL_FAIL,
    _LEVEL_GOOD,
    ArtifactResult,
    JudgeReport,
    _apply_aggregate_verdict,
    _default_brandmind_home,
    _default_output_root,
    _extract_brand_key_text,
    _extract_docx_text,
    _extract_pptx_text,
    _load_transcript_context,
    _quality_level_for_criteria,
)


def _criterion(criterion_id: str, judgment: str) -> dict[str, str]:
    """Build a minimal criterion verdict for quality-level tests."""
    return {
        "id": criterion_id,
        "judgment": judgment,
        "evidence": "evidence",
        "explanation": "explanation",
    }


def _passing_artifact(artifact_type: str) -> ArtifactResult:
    """Build a judged artifact that passes the content gate."""
    return ArtifactResult(
        artifact_type=artifact_type,
        artifact_path=f"/tmp/{artifact_type}",
        passed=True,
        quality_level=_LEVEL_ACCEPTABLE,
    )


def test_required_criteria_met_with_good_gap_is_acceptable() -> None:
    """An artifact is acceptable when required criteria pass but good criteria do not."""
    criteria = [
        _criterion("DOC-1", "MET"),
        _criterion("DOC-2", "MET"),
        _criterion("DOC-3", "MET"),
        _criterion("DOC-4", "MET"),
        _criterion("DOC-5", "UNMET"),
    ]

    level = _quality_level_for_criteria("strategy_document", criteria)

    assert level == _LEVEL_ACCEPTABLE


def test_good_level_requires_good_criterion() -> None:
    """An artifact reaches GOOD only when the artifact-specific good criterion passes."""
    criteria = [
        _criterion("PPT-1", "MET"),
        _criterion("PPT-2", "MET"),
        _criterion("PPT-3", "MET"),
        _criterion("PPT-4", "MET"),
        _criterion("PPT-5", "MET"),
    ]

    level = _quality_level_for_criteria("presentation", criteria)

    assert level == _LEVEL_GOOD


def test_missing_required_criterion_fails_even_with_three_met() -> None:
    """The old three-MET threshold is not enough for strict artifact acceptance."""
    criteria = [
        _criterion("KPI-1", "MET"),
        _criterion("KPI-2", "MET"),
        _criterion("KPI-3", "MET"),
        _criterion("KPI-4", "UNMET"),
        _criterion("KPI-5", "MET"),
        _criterion("KPI-6", "MET"),
        _criterion("KPI-7", "MET"),
    ]

    level = _quality_level_for_criteria("spreadsheet", criteria)

    assert level == _LEVEL_FAIL


def test_aggregate_requires_all_expected_artifacts_to_pass() -> None:
    """A strict session pass requires all expected artifacts to pass."""
    report = JudgeReport(
        session_dir="/tmp/session",
        rubric_path="/tmp/rubric.md",
        judge_model="gemini",
        artifacts=[
            _passing_artifact("brand_key_image"),
            _passing_artifact("strategy_document"),
            _passing_artifact("presentation"),
            _passing_artifact("spreadsheet"),
        ],
    )

    _apply_aggregate_verdict(report)

    assert report.aggregate_pass is True
    assert "4/4 expected artifacts pass" in report.aggregate_summary


def test_aggregate_fails_when_one_artifact_fails() -> None:
    """Three passing artifacts cannot close the strict product gate."""
    report = JudgeReport(
        session_dir="/tmp/session",
        rubric_path="/tmp/rubric.md",
        judge_model="gemini",
        artifacts=[
            _passing_artifact("brand_key_image"),
            _passing_artifact("strategy_document"),
            _passing_artifact("presentation"),
            ArtifactResult(
                artifact_type="spreadsheet",
                artifact_path="/tmp/spreadsheet",
                passed=False,
                quality_level=_LEVEL_FAIL,
            ),
        ],
    )

    _apply_aggregate_verdict(report)

    assert report.aggregate_pass is False
    assert "3/4 expected artifacts pass" in report.aggregate_summary


def test_aggregate_fails_when_artifact_is_skipped() -> None:
    """Skipped artifacts are product gaps and must not lower the pass floor."""
    report = JudgeReport(
        session_dir="/tmp/session",
        rubric_path="/tmp/rubric.md",
        judge_model="gemini",
        artifacts=[
            _passing_artifact("brand_key_image"),
            _passing_artifact("strategy_document"),
            _passing_artifact("presentation"),
            ArtifactResult(
                artifact_type="spreadsheet",
                artifact_path="",
                skipped=True,
                skip_reason="artifact not found on disk",
            ),
        ],
    )

    _apply_aggregate_verdict(report)

    assert report.aggregate_pass is False
    assert "1 missing or skipped" in report.aggregate_summary


def test_pptx_extractor_includes_table_cells(tmp_path) -> None:
    """Presentation judging must see roadmap and KPI tables, not only text boxes."""
    from pptx import Presentation

    path = tmp_path / "deck.pptx"
    deck = Presentation()
    slide = deck.slides.add_slide(deck.slide_layouts[5])
    slide.shapes.title.text = "KPIs & Measurement"
    table = slide.shapes.add_table(2, 2, 914400, 1371600, 7315200, 1371600).table
    table.cell(0, 0).text = "KPI"
    table.cell(0, 1).text = "Target"
    table.cell(1, 0).text = "Weekday bookings"
    table.cell(1, 1).text = "+30% by month 3"
    deck.save(path)

    text, skip_reason = _extract_pptx_text(path)

    assert skip_reason is None
    assert "KPIs & Measurement" in text
    assert "Weekday bookings | +30% by month 3" in text


def test_docx_extractor_includes_table_cells(tmp_path) -> None:
    """Document judging must see roadmap and KPI tables, not only paragraphs."""
    from docx import Document

    path = tmp_path / "strategy.docx"
    document = Document()
    document.add_heading("KPI & Measurement Plan")
    table = document.add_table(rows=2, cols=2)
    table.cell(0, 0).text = "KPI"
    table.cell(0, 1).text = "Target"
    table.cell(1, 0).text = "Weekday bookings"
    table.cell(1, 1).text = "+30% by month 3"
    document.save(path)

    text, skip_reason = _extract_docx_text(path)

    assert skip_reason is None
    assert "KPI & Measurement Plan" in text
    assert "Weekday bookings | +30% by month 3" in text


def test_brand_key_extractor_uses_sidecar_before_ocr(tmp_path) -> None:
    """Brand Key judging should not depend on optional OCR packages."""
    image_path = tmp_path / "brand_key.jpeg"
    image_path.write_bytes(b"not a real image")
    image_path.with_suffix(".brand_key.json").write_text(
        """
        {
          "brand_name": "Chuyện Ba Bữa Signature",
          "components": [
            {"label": "Root Strengths", "value": "Indochine flagship"},
            {"label": "Brand Essence", "value": "Sophisticated Hospitality"}
          ]
        }
        """,
        encoding="utf-8",
    )

    text, skip_reason = _extract_brand_key_text(image_path)

    assert skip_reason is None
    assert "Chuyện Ba Bữa Signature" in text
    assert "Root Strengths: Indochine flagship" in text
    assert "Brand Essence: Sophisticated Hospitality" in text


def test_artifact_judge_transcript_context_accepts_driver_list_schema(tmp_path) -> None:
    """Artifact judge should see API-driver transcripts, not blank turns."""
    session_dir = tmp_path / "pilot"
    session_dir.mkdir()
    (session_dir / "transcript.json").write_text(
        """
        [
          {
            "turn": 1,
            "user_message": "Tôi cần slide pitch.",
            "assistant_response": "Tôi sẽ tạo slide pitch cho founder."
          }
        ]
        """,
        encoding="utf-8",
    )

    text = _load_transcript_context(session_dir)

    assert "USER: Tôi cần slide pitch." in text
    assert "AGENT: Tôi sẽ tạo slide pitch cho founder." in text


def test_artifact_judge_defaults_follow_isolated_run_env(tmp_path, monkeypatch) -> None:
    """Judge audit lookup should follow isolated eval env roots."""
    brandmind_home = tmp_path / "brandmind-home"
    output_root = tmp_path / "brandmind-output"
    monkeypatch.setenv("BRANDMIND_HOME", str(brandmind_home))
    monkeypatch.setenv("BRANDMIND_OUTPUT_DIR", str(output_root))

    assert _default_brandmind_home() == brandmind_home
    assert _default_output_root() == output_root
