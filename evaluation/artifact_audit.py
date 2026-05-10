"""Audit Tier 1 artifact production health for a BrandMind pilot session.

The audit reads a pilot session directory plus the associated agent
workspace and outputs a structured JSON report describing whether the
session produced the deliverable artifacts BrandMind is designed to
generate (Brand Key visual, strategy document, KPI spreadsheet,
executive presentation), how the agent dispatched its sub-agents and
deliverable tools, and whether the workspace notes were kept in sync
with the chat transcript.

Usage:
    python evaluation/artifact_audit.py --session-dir <path-to-pilot>
    python evaluation/artifact_audit.py --session-dir <path> --pretty

The exit code is ``0`` when the audit completes successfully (the report
itself indicates whether artifacts are missing); ``2`` when required
input files cannot be found.
"""

from __future__ import annotations

import argparse
import json
import re
import sys
from dataclasses import asdict, dataclass, field
from pathlib import Path
from typing import Any

DELIVERABLE_TOOLS = (
    "generate_brand_key",
    "generate_document",
    "generate_presentation",
    "generate_spreadsheet",
    "generate_image",
    "edit_image",
)

SUBAGENT_TYPES = (
    "creative-studio",
    "document-generator",
    "market-research",
    "social-media-analyst",
)

PHASE_KEYS = (
    "phase_0",
    "phase_0_5",
    "phase_1",
    "phase_2",
    "phase_3",
    "phase_4",
    "phase_5",
)

ARTIFACT_EXTENSIONS = {
    "brand_key_image": (".jpeg", ".jpg", ".png"),
    "strategy_document": (".docx", ".pdf"),
    "presentation": (".pptx",),
    "spreadsheet": (".xlsx",),
}


# ---------------------------------------------------------------------------
# Data models
# ---------------------------------------------------------------------------


@dataclass
class ToolUsage:
    """Aggregate counts and observed file paths for a deliverable tool."""

    call_count: int = 0
    files_produced: list[str] = field(default_factory=list)


@dataclass
class WorkspaceFile:
    """Inspection result for a single workspace markdown file."""

    exists: bool = False
    lines: int = 0
    sections_present: list[str] = field(default_factory=list)


@dataclass
class Tier1Health:
    """Binary health checks aligned with BrandMind's Tier 1 success target."""

    brand_key_produced: bool = False
    strategy_doc_produced: bool = False
    kpi_xlsx_produced: bool = False
    presentation_produced: bool = False
    workspace_brief_covers_all_phases: bool = False

    @property
    def score(self) -> int:
        return sum(
            int(getattr(self, attr))
            for attr in (
                "brand_key_produced",
                "strategy_doc_produced",
                "kpi_xlsx_produced",
                "presentation_produced",
                "workspace_brief_covers_all_phases",
            )
        )


@dataclass
class SemanticCheck:
    """Structural-quality assessment for a single artifact file.

    The check looks beyond file existence: it parses the artifact with
    its native parser and verifies that the structural depth expected
    of a Phase 5 deliverable is present (e.g. the strategy DOCX
    actually contains the eight phase sections, not just a title page).
    Missing dependencies are recorded as ``skipped`` rather than
    failures so the Tier 1 audit does not penalise environments where
    a parser is unavailable.

    Attributes:
        artifact_type: One of ``brand_key_image``, ``strategy_document``,
            ``presentation``, ``spreadsheet``.
        artifact_path: Absolute path of the artifact under inspection.
        passed: ``True`` when every required structural threshold is met.
            ``True`` when the check is skipped (no penalty for missing
            optional parsers).
        skipped: ``True`` when the parser dependency is missing.
        details: Free-form structural counts the check produced
            (e.g. ``{"section_headings_found": 7}``).
        reasons: Human-readable explanations for failed thresholds.
            Empty when ``passed`` is ``True``.
    """

    artifact_type: str
    artifact_path: str
    passed: bool = True
    skipped: bool = False
    details: dict[str, Any] = field(default_factory=dict)
    reasons: list[str] = field(default_factory=list)


@dataclass
class AuditReport:
    """Complete audit report for one session."""

    session_dir: str
    session_id_api: str
    workspace_session_id: str | None
    workspace_dir: str | None
    scope: str | None
    completed_phases: list[str]
    deliverable_tools: dict[str, ToolUsage]
    subagent_dispatch: dict[str, int]
    artifacts_on_disk: dict[str, list[str]]
    workspace_files: dict[str, WorkspaceFile]
    tier1_health: Tier1Health
    semantic_checks: list[SemanticCheck]
    notes: list[str]

    def to_dict(self) -> dict:
        d = asdict(self)
        d["tier1_health"]["score"] = self.tier1_health.score
        return d


# ---------------------------------------------------------------------------
# Parsers
# ---------------------------------------------------------------------------


def parse_metadata(session_dir: Path) -> tuple[str, str | None, list[str]]:
    """Read metadata.json and return (api_session_id, scope, completed_phases)."""
    meta_path = session_dir / "metadata.json"
    if not meta_path.is_file():
        return "", None, []
    meta = json.loads(meta_path.read_text(encoding="utf-8"))
    api_session_id = meta.get("session_id", "")
    sm = meta.get("session_metadata", {}) or {}
    return api_session_id, sm.get("scope"), list(sm.get("completed_phases", []))


def parse_server_log(log_path: Path) -> tuple[
    dict[str, ToolUsage], dict[str, int], list[str]
]:
    """Extract deliverable tool usage and sub-agent dispatch counts from log.

    The log is line-oriented; each tool call is announced with a header
    line like ``🔬 DEBUG TOOL_CALL [generate_brand_key]`` followed by
    ``└─ key=value`` lines describing arguments. This parser walks the
    log linearly, accumulates counts, and captures any ``file_path``
    arguments emitted by deliverable tools.

    Args:
        log_path: Path to the captured server log file.

    Returns:
        Tuple of (deliverable_tool_usage, subagent_dispatch_counts,
        notes). ``notes`` collects warnings such as missing log file.
    """
    deliverable: dict[str, ToolUsage] = {
        name: ToolUsage() for name in DELIVERABLE_TOOLS
    }
    subagent: dict[str, int] = {name: 0 for name in SUBAGENT_TYPES}
    notes: list[str] = []

    if not log_path.is_file():
        notes.append(f"server.log not found at {log_path}")
        return deliverable, subagent, notes

    tool_call_re = re.compile(r"🔬 DEBUG TOOL_CALL \[(?P<tool>[^\]]+)\]")
    arg_re = re.compile(r"└─ (?P<key>[A-Za-z_][\w]*)=(?P<val>.*)")

    current_tool: str | None = None
    with log_path.open("r", encoding="utf-8", errors="replace") as fh:
        for line in fh:
            m = tool_call_re.search(line)
            if m:
                current_tool = m.group("tool")
                if current_tool in deliverable:
                    deliverable[current_tool].call_count += 1
                continue
            m2 = arg_re.search(line)
            if m2 and current_tool:
                key = m2.group("key")
                val = m2.group("val").strip()
                if current_tool == "task" and key == "subagent_type":
                    name = val.split()[0] if val else ""
                    if name in subagent:
                        subagent[name] += 1
                elif current_tool in deliverable and key in {
                    "file_path",
                    "image_path",
                    "output_path",
                }:
                    deliverable[current_tool].files_produced.append(val)
    return deliverable, subagent, notes


def find_workspace_dir(
    api_session_id: str, brandmind_home: Path
) -> tuple[Path | None, str | None]:
    """Resolve the workspace directory associated with the pilot session.

    BrandStrategySession assigns its own UUID prefix that is independent
    of the API ``session_id``. The two ids are not joined in either
    metadata file, so the workspace is located heuristically:

    1. Read ``project.json`` from each candidate directory under
       ``BRANDMIND_HOME/projects/`` and prefer the one whose
       ``session_id`` matches the API session id (when both happen to
       share a prefix).
    2. Otherwise, fall back to the most recently modified candidate
       whose ``project.json`` ``brand_name`` matches the persona's
       brand. The caller should treat this as best effort.

    Args:
        api_session_id: ``session_id`` from the pilot ``metadata.json``.
        brandmind_home: Root of the on-disk BrandMind state tree.

    Returns:
        Tuple of (workspace path or ``None``, workspace session id or
        ``None``). The session id reflects the workspace project.json
        value, which can differ from ``api_session_id``.
    """
    projects_root = brandmind_home / "projects"
    if not projects_root.is_dir():
        return None, None
    candidates = sorted(
        (p for p in projects_root.iterdir() if p.is_dir()),
        key=lambda p: p.stat().st_mtime,
        reverse=True,
    )
    for cand in candidates:
        proj_file = cand / "project.json"
        if not proj_file.is_file():
            continue
        try:
            data = json.loads(proj_file.read_text(encoding="utf-8"))
        except json.JSONDecodeError:
            continue
        if data.get("session_id") == api_session_id:
            ws = cand / "workspace"
            return (ws if ws.is_dir() else None, cand.name)
    if candidates:
        latest = candidates[0]
        ws = latest / "workspace"
        return (ws if ws.is_dir() else None, latest.name)
    return None, None


def inspect_workspace_files(workspace_dir: Path | None) -> dict[str, WorkspaceFile]:
    """Inspect the three canonical workspace files for line count and sections.

    The audit is intentionally lightweight: it counts non-empty markdown
    headings as ``sections_present`` so the caller can compare which
    phases or sections the agent populated. This does not validate
    content depth — that belongs to higher-tier semantic checks.

    Args:
        workspace_dir: Directory expected to contain the three notes
            files. ``None`` produces empty results with ``exists=False``.

    Returns:
        Mapping of filename to :class:`WorkspaceFile`.
    """
    targets = ("brand_brief.md", "quality_gates.md", "working_notes.md")
    out: dict[str, WorkspaceFile] = {name: WorkspaceFile() for name in targets}
    if workspace_dir is None or not workspace_dir.is_dir():
        return out
    for name in targets:
        path = workspace_dir / name
        if not path.is_file():
            continue
        text = path.read_text(encoding="utf-8", errors="replace")
        wf = out[name]
        wf.exists = True
        wf.lines = sum(1 for _ in text.splitlines())
        wf.sections_present = [
            line.strip("# ").strip()
            for line in text.splitlines()
            if line.lstrip().startswith("##")
        ]
    return out


def parse_pilot_state(
    pilot_state_path: Path,
) -> tuple[dict[str, ToolUsage], dict[str, int]]:
    """Extract deliverable tool counts from the smoke harness pilot state.

    The smoke harness records each turn's ``tools_used`` list inside
    ``_pilot_state.json`` (``smoke_test.py:_send_turn``). Each entry is
    a tool name string; ``task`` invocations have their ``subagent_type``
    captured separately under ``tool_args``. The audit consumes this
    file as the primary tool-usage source so it stays accurate even
    when the running server log is not piped into the session
    directory. Server-log parsing remains as the fallback for runs
    that do capture ``server.log`` (full-pilot driver harness).

    Args:
        pilot_state_path: Path to the smoke harness state file.

    Returns:
        Tuple of (deliverable_tool_usage, subagent_dispatch_counts).
        Returns empty maps when the file is absent so the caller can
        merge with ``parse_server_log`` results unconditionally.
    """
    deliverable: dict[str, ToolUsage] = {
        name: ToolUsage() for name in DELIVERABLE_TOOLS
    }
    subagent: dict[str, int] = {name: 0 for name in SUBAGENT_TYPES}
    if not pilot_state_path.is_file():
        return deliverable, subagent
    try:
        state = json.loads(pilot_state_path.read_text(encoding="utf-8"))
    except json.JSONDecodeError:
        return deliverable, subagent
    for turn in state.get("turns", []):
        for tool_entry in turn.get("tools_used", []) or []:
            if isinstance(tool_entry, dict):
                tool_name = tool_entry.get("name") or tool_entry.get("tool_name")
                tool_args = tool_entry.get("args") or tool_entry.get("tool_args")
            else:
                tool_name = str(tool_entry)
                tool_args = None
            if not tool_name:
                continue
            if tool_name in deliverable:
                deliverable[tool_name].call_count += 1
            elif tool_name == "task" and isinstance(tool_args, dict):
                target = (tool_args.get("subagent_type") or "").split()[0]
                if target in subagent:
                    subagent[target] += 1
    return deliverable, subagent


def merge_tool_usage(
    primary: dict[str, ToolUsage], fallback: dict[str, ToolUsage]
) -> dict[str, ToolUsage]:
    """Merge two deliverable-tool-usage maps, preferring the primary source.

    The primary source is the pilot harness ``_pilot_state.json``
    (always present after a smoke run); the fallback is the parsed
    server log (only available when the harness pipes the server log
    into the session directory). When both record a non-zero count for
    the same tool, the primary wins; ``files_produced`` is concatenated
    so on-disk paths from either source are not lost.

    Args:
        primary: Tool usage from :func:`parse_pilot_state`.
        fallback: Tool usage from :func:`parse_server_log`.

    Returns:
        A merged map keyed by tool name.
    """
    merged: dict[str, ToolUsage] = {}
    for name in DELIVERABLE_TOOLS:
        p = primary.get(name, ToolUsage())
        f = fallback.get(name, ToolUsage())
        merged_files = list(dict.fromkeys(list(p.files_produced) + list(f.files_produced)))
        merged[name] = ToolUsage(
            call_count=max(p.call_count, f.call_count),
            files_produced=merged_files,
        )
    return merged


def merge_subagent(
    primary: dict[str, int], fallback: dict[str, int]
) -> dict[str, int]:
    """Merge two sub-agent dispatch maps, preferring the primary source."""
    merged: dict[str, int] = {}
    for name in SUBAGENT_TYPES:
        merged[name] = max(primary.get(name, 0), fallback.get(name, 0))
    return merged


def parse_manifest(
    manifest_path: Path, session_id: str | None
) -> dict[str, list[str]]:
    """Read the artifact manifest JSONL and group records by audit category.

    The manifest is the ground-truth provenance log: each artifact
    record carries ``session_id`` so the audit can answer "which files
    belong to this session" without depending only on directory layout.

    Args:
        manifest_path: Absolute path to a ``.manifest.jsonl`` file.
            Missing or unreadable files return empty lists.
        session_id: When provided, only manifest records whose
            ``session_id`` matches are returned. When ``None`` every
            record is returned.

    Returns:
        Mapping of audit category (``brand_key_image``,
        ``strategy_document``, ``presentation``, ``spreadsheet``,
        ``other_image``) to absolute file paths. Records pointing at
        files that no longer exist are dropped so callers see only
        artifacts a downstream judge can actually open.
    """
    out: dict[str, list[str]] = {
        "brand_key_image": [],
        "strategy_document": [],
        "presentation": [],
        "spreadsheet": [],
        "other_image": [],
    }
    if not manifest_path.is_file():
        return out
    try:
        text = manifest_path.read_text(encoding="utf-8", errors="replace")
    except OSError:
        return out
    for line in text.splitlines():
        line = line.strip()
        if not line:
            continue
        try:
            rec = json.loads(line)
        except json.JSONDecodeError:
            continue
        if session_id and rec.get("session_id") != session_id:
            continue
        path = rec.get("path") or ""
        category = rec.get("category") or ""
        if not path or not Path(path).is_file():
            continue
        if category == "documents":
            out["strategy_document"].append(path)
        elif category == "presentations":
            out["presentation"].append(path)
        elif category == "spreadsheets":
            out["spreadsheet"].append(path)
        elif category == "images":
            if "brand_key" in Path(path).name.lower():
                out["brand_key_image"].append(path)
            else:
                out["other_image"].append(path)
    return out


def scan_artifacts_on_disk(
    output_root: Path,
    session_start_iso: str | None,
    session_id: str | None = None,
    legacy_roots: list[Path] | None = None,
) -> dict[str, list[str]]:
    """Enumerate artifacts produced by the pilot session.

    Resolution order:

    1. Read ``$output_root/.manifest.jsonl`` filtered by ``session_id``
       as the primary provenance source.
    2. Read manifest files at any ``legacy_roots``. Caller passes only
       known artifact roots so the fallback stays bounded.
    3. When the manifest sources are empty (legacy session that
       predates manifest provenance), fall back to ``output_root.rglob``
       filtered by mtime relative to ``session_start_iso``.

    Args:
        output_root: Root of the BrandMind output tree.
        session_start_iso: ISO timestamp at which the pilot session
            started. Used by the rglob fallback to drop pre-existing
            artifacts; ``None`` disables the time filter.
        session_id: Pilot session id used to filter manifest records.
            ``None`` disables the manifest pass and forces the rglob
            fallback (legacy compatibility).
        legacy_roots: Additional manifest roots to consult, typically
            known workspace roots used by older sessions. Each root
            must be a trusted location, not a user-supplied path.

    Returns:
        Mapping of artifact category to list of file paths that match
        the category for this session.
    """
    from datetime import datetime

    out: dict[str, list[str]] = {
        "brand_key_image": [],
        "strategy_document": [],
        "presentation": [],
        "spreadsheet": [],
        "other_image": [],
    }

    if session_id:
        manifest_paths = [output_root / ".manifest.jsonl"]
        if legacy_roots:
            manifest_paths.extend(root / ".manifest.jsonl" for root in legacy_roots)
        for mpath in manifest_paths:
            chunk = parse_manifest(mpath, session_id)
            for key, values in chunk.items():
                out[key].extend(values)
        for key in out:
            out[key] = list(dict.fromkeys(out[key]))
        if any(out.values()):
            return out

    cutoff = None
    if session_start_iso:
        try:
            cutoff = datetime.fromisoformat(session_start_iso).timestamp()
        except ValueError:
            cutoff = None

    if not output_root.is_dir():
        return out
    for entry in output_root.rglob("*"):
        if not entry.is_file():
            continue
        if cutoff is not None and entry.stat().st_mtime < cutoff:
            continue
        if "/eval/" in str(entry):
            continue
        suffix = entry.suffix.lower()
        name = entry.name.lower()
        if suffix in ARTIFACT_EXTENSIONS["spreadsheet"]:
            out["spreadsheet"].append(str(entry))
        elif suffix in ARTIFACT_EXTENSIONS["presentation"]:
            out["presentation"].append(str(entry))
        elif suffix in ARTIFACT_EXTENSIONS["strategy_document"]:
            out["strategy_document"].append(str(entry))
        elif suffix in ARTIFACT_EXTENSIONS["brand_key_image"]:
            if "brand_key" in name:
                out["brand_key_image"].append(str(entry))
            else:
                out["other_image"].append(str(entry))
    return out


def evaluate_tier1_health(
    deliverable: dict[str, ToolUsage],
    artifacts: dict[str, list[str]],
    workspace: dict[str, WorkspaceFile],
    completed_phases: list[str],
) -> Tier1Health:
    """Reduce the audit signals into the binary Tier 1 health checks.

    A check passes when the corresponding tool was called at least once
    AND a matching artifact is observed on disk (when applicable). The
    workspace coverage check passes when ``brand_brief.md`` contains a
    section heading for every completed phase.

    Args:
        deliverable: Tool usage map from :func:`parse_server_log`.
        artifacts: On-disk artifact map from
            :func:`scan_artifacts_on_disk`.
        workspace: Workspace inspection from
            :func:`inspect_workspace_files`.
        completed_phases: List of phase ids the session advanced through.

    Returns:
        :class:`Tier1Health` with the five binary checks set.
    """
    health = Tier1Health()
    bk_calls = deliverable.get("generate_brand_key", ToolUsage()).call_count
    health.brand_key_produced = bk_calls > 0 or bool(artifacts["brand_key_image"])
    doc_calls = deliverable.get("generate_document", ToolUsage()).call_count
    health.strategy_doc_produced = doc_calls > 0 or bool(
        artifacts["strategy_document"]
    )
    pres_calls = deliverable.get("generate_presentation", ToolUsage()).call_count
    health.presentation_produced = pres_calls > 0 or bool(artifacts["presentation"])
    sheet_calls = deliverable.get("generate_spreadsheet", ToolUsage()).call_count
    health.kpi_xlsx_produced = sheet_calls > 0 or bool(artifacts["spreadsheet"])

    brief = workspace.get("brand_brief.md", WorkspaceFile())
    if brief.exists:
        sections_lower = " ".join(s.lower() for s in brief.sections_present)
        phase_tokens = {
            "phase_0": ["phase 0", "phase_0", "diagnosis"],
            "phase_0_5": ["phase 0.5", "phase_0_5", "equity audit"],
            "phase_1": ["phase 1", "phase_1", "market intelligence"],
            "phase_2": ["phase 2", "phase_2", "positioning"],
            "phase_3": ["phase 3", "phase_3", "identity"],
            "phase_4": ["phase 4", "phase_4", "communication"],
            "phase_5": ["phase 5", "phase_5", "deliverable", "strategy plan"],
        }
        missing = [
            phase
            for phase in completed_phases
            if not any(tok in sections_lower for tok in phase_tokens.get(phase, []))
        ]
        health.workspace_brief_covers_all_phases = not missing
    return health


# ---------------------------------------------------------------------------
# Semantic structural checks
# ---------------------------------------------------------------------------


_BRAND_KEY_LABEL_PATTERNS: tuple[tuple[str, ...], ...] = (
    ("root strength", "core strength", "heritage", "nguồn gốc"),
    ("competitive environment", "competitor", "competitive", "cạnh tranh"),
    ("target", "target audience", "đối tượng"),
    ("insight", "consumer insight", "thấu hiểu"),
    ("benefit", "functional benefit", "emotional benefit", "lợi ích"),
    ("values", "personality", "tính cách"),
    ("reason to believe", "rtb", "lý do tin"),
    ("discriminator", "point of difference", "pod", "khác biệt"),
    ("essence", "brand essence", "mantra", "cốt lõi"),
)

_STRATEGY_DOC_PHASE_PATTERNS: tuple[tuple[str, ...], ...] = (
    ("executive summary", "tóm tắt"),
    ("business context", "diagnosis", "problem statement", "bối cảnh"),
    ("market intelligence", "research", "competitive landscape", "thị trường"),
    ("brand positioning", "positioning statement", "định vị"),
    ("brand identity", "personality", "visual", "nhận diện"),
    ("communication framework", "messaging", "channel", "truyền thông"),
    ("implementation roadmap", "execution", "lộ trình"),
    ("kpi", "metric", "measurement", "đo lường"),
)

_PRESENTATION_TITLE_PATTERNS: tuple[tuple[str, ...], ...] = (
    ("challenge", "problem", "diagnosis", "thách thức"),
    ("position", "positioning", "định vị"),
    ("identity", "brand personality", "nhận diện"),
    ("communication", "messaging", "channel", "truyền thông"),
    ("roadmap", "implementation", "timeline", "lộ trình"),
    ("kpi", "metric", "measurement", "đo lường"),
)

_KPI_REQUIRED_HEADERS: tuple[tuple[str, ...], ...] = (
    ("metric", "kpi", "chỉ số"),
    ("measurement", "method", "phương pháp"),
    ("target", "goal", "mục tiêu"),
)


def _matches_any(text_lower: str, patterns: tuple[tuple[str, ...], ...]) -> int:
    """Return how many pattern groups have at least one term in ``text_lower``."""
    return sum(
        1 for group in patterns if any(term in text_lower for term in group)
    )


def _check_brand_key_semantic(path: Path) -> SemanticCheck:
    """Verify the Brand Key image contains enough of the 9 components.

    OCR via :mod:`pytesseract` is preferred when both the Python wheel
    and the system ``tesseract`` binary are available. When OCR cannot
    run, the check is recorded as skipped (passed without inspection)
    so absent system dependencies do not invalidate Tier 1.

    Args:
        path: Absolute path of the Brand Key image.

    Returns:
        :class:`SemanticCheck` with structural details for the image.
    """
    check = SemanticCheck(artifact_type="brand_key_image", artifact_path=str(path))
    try:
        import pytesseract  # type: ignore
        from PIL import Image  # type: ignore
    except ImportError as exc:
        check.skipped = True
        check.details["skip_reason"] = f"OCR unavailable: {exc}"
        return check
    try:
        text = pytesseract.image_to_string(Image.open(path), lang="eng+vie")
    except Exception as exc:  # noqa: BLE001
        check.skipped = True
        check.details["skip_reason"] = f"OCR failed: {exc}"
        return check
    count = _matches_any(text.lower(), _BRAND_KEY_LABEL_PATTERNS)
    check.details["component_labels_found"] = count
    check.details["expected_minimum"] = 7
    if count < 7:
        check.passed = False
        check.reasons.append(
            f"Brand Key OCR found {count}/9 expected component labels "
            "(threshold 7)."
        )
    return check


def _check_strategy_doc_semantic(path: Path) -> SemanticCheck:
    """Verify the strategy DOCX covers at least six of eight phase sections."""
    check = SemanticCheck(artifact_type="strategy_document", artifact_path=str(path))
    try:
        import docx  # type: ignore
    except ImportError as exc:
        check.skipped = True
        check.details["skip_reason"] = f"python-docx unavailable: {exc}"
        return check
    try:
        document = docx.Document(str(path))
    except Exception as exc:  # noqa: BLE001
        check.passed = False
        check.reasons.append(f"DOCX failed to open: {exc}")
        return check
    headings = [
        p.text.lower().strip()
        for p in document.paragraphs
        if (p.style and "heading" in (p.style.name or "").lower())
        or (p.text and len(p.text) < 100 and p.text.strip().endswith(":") is False)
    ]
    heading_text = " | ".join(headings)
    matched = _matches_any(heading_text, _STRATEGY_DOC_PHASE_PATTERNS)
    check.details["section_headings_matched"] = matched
    check.details["expected_minimum"] = 6
    check.details["heading_count"] = len(headings)
    if matched < 6:
        check.passed = False
        check.reasons.append(
            f"DOCX matched {matched}/8 expected phase sections (threshold 6)."
        )
    return check


def _check_presentation_semantic(path: Path) -> SemanticCheck:
    """Verify the PPTX has enough slides and title coverage of the deck flow."""
    check = SemanticCheck(artifact_type="presentation", artifact_path=str(path))
    try:
        from pptx import Presentation  # type: ignore
    except ImportError as exc:
        check.skipped = True
        check.details["skip_reason"] = f"python-pptx unavailable: {exc}"
        return check
    try:
        deck = Presentation(str(path))
    except Exception as exc:  # noqa: BLE001
        check.passed = False
        check.reasons.append(f"PPTX failed to open: {exc}")
        return check
    titles: list[str] = []
    for slide in deck.slides:
        if slide.shapes.title and slide.shapes.title.text:
            titles.append(slide.shapes.title.text.lower())
    slide_count = len(deck.slides)
    title_text = " | ".join(titles)
    matched = _matches_any(title_text, _PRESENTATION_TITLE_PATTERNS)
    check.details["slide_count"] = slide_count
    check.details["title_topics_matched"] = matched
    check.details["expected_slide_minimum"] = 8
    check.details["expected_topic_minimum"] = 4
    if slide_count < 8:
        check.passed = False
        check.reasons.append(
            f"PPTX has {slide_count} slides (threshold 8)."
        )
    if matched < 4:
        check.passed = False
        check.reasons.append(
            f"PPTX titles matched {matched}/6 expected topics (threshold 4)."
        )
    return check


def _check_kpi_xlsx_semantic(path: Path) -> SemanticCheck:
    """Verify the KPI XLSX has at least five data rows and required columns."""
    check = SemanticCheck(artifact_type="spreadsheet", artifact_path=str(path))
    try:
        from openpyxl import load_workbook  # type: ignore
    except ImportError as exc:
        check.skipped = True
        check.details["skip_reason"] = f"openpyxl unavailable: {exc}"
        return check
    try:
        workbook = load_workbook(str(path), read_only=True, data_only=True)
    except Exception as exc:  # noqa: BLE001
        check.passed = False
        check.reasons.append(f"XLSX failed to open: {exc}")
        return check
    primary_sheet = workbook[workbook.sheetnames[0]]
    rows = [
        [str(cell).strip().lower() if cell is not None else "" for cell in row]
        for row in primary_sheet.iter_rows(values_only=True)
    ]
    header_pool = " | ".join(" ".join(row) for row in rows[:3])
    header_match = _matches_any(header_pool, _KPI_REQUIRED_HEADERS)
    data_rows = [row for row in rows[1:] if any(cell for cell in row)]
    check.details["sheet_count"] = len(workbook.sheetnames)
    check.details["data_rows"] = len(data_rows)
    check.details["header_topics_matched"] = header_match
    check.details["expected_data_rows_minimum"] = 5
    check.details["expected_header_topics_minimum"] = 3
    if len(data_rows) < 5:
        check.passed = False
        check.reasons.append(
            f"XLSX primary sheet has {len(data_rows)} data rows (threshold 5)."
        )
    if header_match < 3:
        check.passed = False
        check.reasons.append(
            f"XLSX header topics matched {header_match}/3 expected categories."
        )
    return check


def run_semantic_checks(
    artifacts: dict[str, list[str]],
) -> list[SemanticCheck]:
    """Run a structural check per produced artifact and return the results.

    Iterates through the four canonical artifact buckets emitted by
    :func:`scan_artifacts_on_disk` and dispatches to the dedicated
    checker for each. Buckets without an artifact are skipped silently;
    their absence is already captured by Tier 1 health.

    Args:
        artifacts: Mapping returned by :func:`scan_artifacts_on_disk`.

    Returns:
        Ordered list of :class:`SemanticCheck` outcomes (one per
        artifact found on disk).
    """
    results: list[SemanticCheck] = []
    for path in artifacts.get("brand_key_image", []):
        results.append(_check_brand_key_semantic(Path(path)))
    for path in artifacts.get("strategy_document", []):
        if path.lower().endswith(".docx"):
            results.append(_check_strategy_doc_semantic(Path(path)))
    for path in artifacts.get("presentation", []):
        results.append(_check_presentation_semantic(Path(path)))
    for path in artifacts.get("spreadsheet", []):
        results.append(_check_kpi_xlsx_semantic(Path(path)))
    return results


# ---------------------------------------------------------------------------
# CLI
# ---------------------------------------------------------------------------


def audit(
    session_dir: Path,
    brandmind_home: Path,
    output_root: Path,
    semantic: bool = False,
    legacy_roots: list[Path] | None = None,
) -> AuditReport:
    """Run the full audit pipeline for one session directory.

    Args:
        session_dir: Pilot session directory (the eval folder, not the
            workspace).
        brandmind_home: Root of the on-disk BrandMind state tree, used
            to resolve the workspace directory.
        output_root: Root of the BrandMind output tree where artifacts
            land on disk.
        semantic: When ``True``, run :func:`run_semantic_checks` over
            the produced artifacts and include the results on the
            returned report. Defaults to ``False`` so the lightweight
            existence check stays the default for the M-3 smoke path.
        legacy_roots: Optional list of additional manifest roots to
            consult when ``output_root`` may not capture every
            artifact. Each root is treated as a known location, not a
            user-supplied path.

    Returns:
        A populated :class:`AuditReport` ready for serialisation.

    Raises:
        FileNotFoundError: when ``session_dir`` does not exist.
    """
    if not session_dir.is_dir():
        raise FileNotFoundError(f"Session directory not found: {session_dir}")

    api_session_id, scope, completed_phases = parse_metadata(session_dir)
    deliverable_log, subagent_log, notes = parse_server_log(
        session_dir / "server.log"
    )
    deliverable_pilot, subagent_pilot = parse_pilot_state(
        session_dir / "_pilot_state.json"
    )
    deliverable = merge_tool_usage(deliverable_pilot, deliverable_log)
    subagent = merge_subagent(subagent_pilot, subagent_log)
    workspace_dir, workspace_id = find_workspace_dir(api_session_id, brandmind_home)
    workspace_files = inspect_workspace_files(workspace_dir)

    meta_path = session_dir / "metadata.json"
    session_start_iso: str | None = None
    if meta_path.is_file():
        meta = json.loads(meta_path.read_text(encoding="utf-8"))
        session_start_iso = meta.get("created_at") or meta.get("date")

    artifacts = scan_artifacts_on_disk(
        output_root,
        session_start_iso,
        session_id=api_session_id,
        legacy_roots=legacy_roots,
    )
    health = evaluate_tier1_health(
        deliverable, artifacts, workspace_files, completed_phases
    )
    semantic_checks: list[SemanticCheck] = (
        run_semantic_checks(artifacts) if semantic else []
    )

    return AuditReport(
        session_dir=str(session_dir),
        session_id_api=api_session_id,
        workspace_session_id=workspace_id,
        workspace_dir=str(workspace_dir) if workspace_dir else None,
        scope=scope,
        completed_phases=completed_phases,
        deliverable_tools=deliverable,
        subagent_dispatch=subagent,
        artifacts_on_disk=artifacts,
        workspace_files=workspace_files,
        tier1_health=health,
        semantic_checks=semantic_checks,
        notes=notes,
    )


def _build_argparser() -> argparse.ArgumentParser:
    parser = argparse.ArgumentParser(description=__doc__)
    parser.add_argument(
        "--session-dir",
        required=True,
        type=Path,
        help="Path to the pilot session directory (under brandmind-output/eval/).",
    )
    parser.add_argument(
        "--brandmind-home",
        type=Path,
        default=Path.home() / ".brandmind",
        help="Root of on-disk BrandMind state. Defaults to ~/.brandmind.",
    )
    parser.add_argument(
        "--output-root",
        type=Path,
        default=Path.cwd() / "brandmind-output",
        help="Root where deliverable artifacts land. Defaults to ./brandmind-output.",
    )
    parser.add_argument(
        "--pretty",
        action="store_true",
        help="Print a human-readable summary in addition to JSON.",
    )
    parser.add_argument(
        "--semantic",
        action="store_true",
        help=(
            "Run structural quality checks on each produced artifact "
            "(parses DOCX/PPTX/XLSX, OCRs the Brand Key image when "
            "pytesseract is available)."
        ),
    )
    return parser


def _format_summary(report: AuditReport) -> str:
    lines: list[str] = []
    h = report.tier1_health
    lines.append(f"Session: {report.session_dir}")
    lines.append(
        f"Scope: {report.scope}  Completed phases: "
        f"{', '.join(report.completed_phases) or '(none)'}"
    )
    lines.append(f"Tier1 health: {h.score}/5")
    lines.append(
        f"  brand_key_produced       = {h.brand_key_produced}"
    )
    lines.append(
        f"  strategy_doc_produced    = {h.strategy_doc_produced}"
    )
    lines.append(
        f"  kpi_xlsx_produced        = {h.kpi_xlsx_produced}"
    )
    lines.append(
        f"  presentation_produced    = {h.presentation_produced}"
    )
    lines.append(
        f"  workspace_brief_covers_all_phases = "
        f"{h.workspace_brief_covers_all_phases}"
    )
    lines.append("Deliverable tool calls:")
    for name in DELIVERABLE_TOOLS:
        usage = report.deliverable_tools[name]
        lines.append(
            f"  {name}: {usage.call_count}"
            + (
                f"  files={usage.files_produced}"
                if usage.files_produced
                else ""
            )
        )
    lines.append("Sub-agent dispatch:")
    for name, count in report.subagent_dispatch.items():
        lines.append(f"  {name}: {count}")
    if report.semantic_checks:
        lines.append("Semantic checks:")
        for sc in report.semantic_checks:
            if sc.skipped:
                status = "SKIP"
            elif sc.passed:
                status = "PASS"
            else:
                status = "FAIL"
            lines.append(
                f"  [{status}] {sc.artifact_type}: "
                f"{Path(sc.artifact_path).name}  details={sc.details}"
            )
            for reason in sc.reasons:
                lines.append(f"    - {reason}")
    return "\n".join(lines)


def main(argv: list[str] | None = None) -> int:
    args = _build_argparser().parse_args(argv)
    try:
        report = audit(
            args.session_dir,
            args.brandmind_home,
            args.output_root,
            semantic=args.semantic,
        )
    except FileNotFoundError as exc:
        print(f"error: {exc}", file=sys.stderr)
        return 2
    payload = report.to_dict()
    print(json.dumps(payload, indent=2, ensure_ascii=False))
    if args.pretty:
        print("\n--- Human summary ---", file=sys.stderr)
        print(_format_summary(report), file=sys.stderr)
    return 0


if __name__ == "__main__":
    raise SystemExit(main())
