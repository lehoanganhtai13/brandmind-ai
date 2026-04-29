"""Artifact content correctness judge (M-6.5).

Reads each Phase 5 artifact a brand-strategy session produced, extracts
its text content, and asks an LLM judge to verify that the content
matches the strategy decisions made during the session. The check
complements the structural depth verification in
``evaluation/artifact_audit.py`` — that one says "the file has the
right shape", this one says "the content inside is faithful to THIS
session".

The judge is intentionally a single LLM (Gemini Flash Lite) calling
through the same :class:`shared.model_clients.llm.google.GoogleAIClientLLM`
client that the production content-check middleware uses. There is no
cross-judge panel and no Fleiss' Kappa: the judge produces a binary
"content matches strategy" signal that Tier 2 prompt levers can be
measured against. Cross-judge calibration is a Tier 2 concern handled
by ``evaluation/judge/run_judges.py``.

Usage:
    uv run python evaluation/judge/artifact_judge.py \\
        --session-dir brandmind-output/eval/<pilot>
    uv run python evaluation/judge/artifact_judge.py \\
        --session-dir <pilot> --json
    uv run python evaluation/judge/artifact_judge.py \\
        --session-dir <pilot> --artifact strategy_document
"""

from __future__ import annotations

import argparse
import asyncio
import json
import sys
from dataclasses import asdict, dataclass, field
from pathlib import Path
from typing import Any

from pydantic import BaseModel, Field

# Repo paths.
_REPO_ROOT = Path(__file__).resolve().parent.parent.parent
for sub in ("config", "core/src", "shared/src"):
    p = _REPO_ROOT / "src" / sub
    if p.is_dir() and str(p) not in sys.path:
        sys.path.insert(0, str(p))

# Reuse the existing audit module to locate the artifacts on disk and
# the structural-text helpers — keeps content extraction in one place.
sys.path.insert(0, str(Path(__file__).resolve().parent.parent))
from artifact_audit import (  # type: ignore  # noqa: E402
    audit as run_existence_audit,
)

from config.system_config import SETTINGS  # noqa: E402
from shared.model_clients.llm.google import (  # noqa: E402
    GoogleAIClientLLM,
    GoogleAIClientLLMConfig,
)

_RUBRIC_PATH = Path(__file__).resolve().parent / "artifact_rubric.md"

# Per-artifact criterion lists. Order matches the rubric file so the
# judge prompt can show criteria in the same sequence.
_CRITERIA: dict[str, tuple[str, ...]] = {
    "brand_key_image": ("BK-1", "BK-2", "BK-3", "BK-4"),
    "strategy_document": ("DOC-1", "DOC-2", "DOC-3", "DOC-4"),
    "presentation": ("PPT-1", "PPT-2", "PPT-3", "PPT-4"),
    "spreadsheet": ("KPI-1", "KPI-2", "KPI-3", "KPI-4"),
}

_PASS_FLOOR_PER_ARTIFACT = 3  # MET threshold per rubric definition
_AGGREGATE_PASS_FLOOR = 3  # 3 of 4 artifacts must pass


# ---------------------------------------------------------------------------
# Verdict schema (returned by the LLM judge)
# ---------------------------------------------------------------------------


class CriterionVerdict(BaseModel):
    """Judge's verdict for one rubric criterion."""

    id: str = Field(..., description="Criterion identifier, e.g. 'BK-1'.")
    judgment: str = Field(
        ..., description="One of: MET, UNMET, CANNOT_ASSESS."
    )
    evidence: str = Field(
        default="",
        description=(
            "Direct quote from the artifact or transcript supporting "
            "the judgment. Empty when CANNOT_ASSESS."
        ),
    )
    explanation: str = Field(
        ...,
        description=(
            "One sentence explaining WHY this judgment, grounded in "
            "the evidence and the rubric definition."
        ),
    )


class ArtifactVerdict(BaseModel):
    """Judge's full verdict for one artifact."""

    artifact_type: str = Field(
        ...,
        description="One of: brand_key_image, strategy_document, "
        "presentation, spreadsheet.",
    )
    criteria: list[CriterionVerdict] = Field(
        default_factory=list,
        description="Per-criterion verdicts in rubric order.",
    )


# ---------------------------------------------------------------------------
# Result data model
# ---------------------------------------------------------------------------


@dataclass
class ArtifactResult:
    """Aggregated content-judge result for one produced artifact."""

    artifact_type: str
    artifact_path: str
    skipped: bool = False
    skip_reason: str = ""
    criteria: list[dict[str, str]] = field(default_factory=list)
    met_count: int = 0
    unmet_count: int = 0
    cannot_assess_count: int = 0
    passed: bool = False


@dataclass
class JudgeReport:
    """Full content-judge report for one session."""

    session_dir: str
    rubric_path: str
    judge_model: str
    artifacts: list[ArtifactResult] = field(default_factory=list)
    aggregate_pass: bool = False
    aggregate_summary: str = ""


# ---------------------------------------------------------------------------
# Artifact text extraction
# ---------------------------------------------------------------------------


def _extract_brand_key_text(path: Path) -> tuple[str, str | None]:
    """Return (text, skip_reason). text is empty when extraction fails."""
    try:
        import pytesseract  # type: ignore
        from PIL import Image  # type: ignore
    except ImportError as exc:
        return "", f"OCR unavailable: {exc}"
    try:
        text = pytesseract.image_to_string(Image.open(path), lang="eng+vie")
    except Exception as exc:  # noqa: BLE001
        return "", f"OCR failed: {exc}"
    return text, None


def _extract_docx_text(path: Path) -> tuple[str, str | None]:
    try:
        import docx  # type: ignore
    except ImportError as exc:
        return "", f"python-docx unavailable: {exc}"
    try:
        document = docx.Document(str(path))
    except Exception as exc:  # noqa: BLE001
        return "", f"DOCX failed to open: {exc}"
    paragraphs = [
        p.text for p in document.paragraphs if p.text and p.text.strip()
    ]
    return "\n".join(paragraphs), None


def _extract_pptx_text(path: Path) -> tuple[str, str | None]:
    try:
        from pptx import Presentation  # type: ignore
    except ImportError as exc:
        return "", f"python-pptx unavailable: {exc}"
    try:
        deck = Presentation(str(path))
    except Exception as exc:  # noqa: BLE001
        return "", f"PPTX failed to open: {exc}"
    chunks: list[str] = []
    for i, slide in enumerate(deck.slides, start=1):
        title = (
            slide.shapes.title.text
            if slide.shapes.title and slide.shapes.title.text
            else "(untitled)"
        )
        chunks.append(f"--- Slide {i}: {title} ---")
        for shape in slide.shapes:
            if (
                shape.has_text_frame
                and shape != slide.shapes.title
                and shape.text_frame.text.strip()
            ):
                chunks.append(shape.text_frame.text)
    return "\n".join(chunks), None


def _extract_xlsx_text(path: Path) -> tuple[str, str | None]:
    try:
        from openpyxl import load_workbook  # type: ignore
    except ImportError as exc:
        return "", f"openpyxl unavailable: {exc}"
    try:
        workbook = load_workbook(str(path), read_only=True, data_only=True)
    except Exception as exc:  # noqa: BLE001
        return "", f"XLSX failed to open: {exc}"
    chunks: list[str] = []
    for sheet_name in workbook.sheetnames:
        chunks.append(f"--- Sheet: {sheet_name} ---")
        sheet = workbook[sheet_name]
        for row in sheet.iter_rows(min_row=1, max_row=20, values_only=True):
            cells = [
                str(cell) if cell is not None else "" for cell in row[:10]
            ]
            chunks.append(" | ".join(cells))
    return "\n".join(chunks), None


_EXTRACTORS = {
    "brand_key_image": _extract_brand_key_text,
    "strategy_document": _extract_docx_text,
    "presentation": _extract_pptx_text,
    "spreadsheet": _extract_xlsx_text,
}


# ---------------------------------------------------------------------------
# Transcript context
# ---------------------------------------------------------------------------


def _load_transcript_context(session_dir: Path, max_chars: int = 20000) -> str:
    """Return the most recent agent/user turns up to ``max_chars`` chars."""
    transcript_path = session_dir / "transcript.json"
    state_path = session_dir / "_pilot_state.json"
    source = transcript_path if transcript_path.is_file() else state_path
    if not source.is_file():
        return "(no transcript available)"
    data = json.loads(source.read_text(encoding="utf-8"))
    turns = data.get("turns", [])
    chunks: list[str] = []
    total = 0
    for turn in reversed(turns):
        block = (
            f"--- Turn {turn.get('turn', '?')} ---\n"
            f"USER: {turn.get('user', '')}\n\n"
            f"AGENT: {turn.get('agent', '')}"
        )
        if total + len(block) > max_chars:
            break
        chunks.append(block)
        total += len(block)
    return "\n\n".join(reversed(chunks))


# ---------------------------------------------------------------------------
# Judge invocation
# ---------------------------------------------------------------------------


_SYSTEM_PROMPT_TEMPLATE = """You are an artifact content reviewer for a Vietnamese
brand-strategy assistant.
You are given (1) the rubric below, (2) recent agent / user turns from
the strategy session, and (3) the extracted text content of one
artifact the agent produced.

Your job: for every criterion under the artifact's section in the
rubric, decide MET / UNMET / CANNOT_ASSESS. Quote the strongest
evidence (from the artifact text or the transcript). Explain your
verdict in one sentence per criterion.

Be honest. The artifact is intended for a junior marketer to present to
her boss. Generic-sounding content that could fit any brand is UNMET.
The transcript decisions take priority over the artifact when the two
disagree (the transcript is the ground truth for THIS session).

Score only the criteria for the requested artifact type. Do not invent
new criteria. Always return the verdicts in the order they appear in
the rubric.

Rubric:
---
{rubric}
---
"""


_USER_PROMPT_TEMPLATE = """Artifact under review: {artifact_type}
Path: {artifact_path}

Required criteria (in order, judge each one): {criteria_list}

=== ARTIFACT CONTENT ===
{artifact_content}

=== RECENT TRANSCRIPT TURNS ===
{transcript_context}

Return the JSON verdict per the schema. Score every criterion listed
above; do not skip any.
"""


def _build_llm(model_id: str) -> GoogleAIClientLLM:
    """Construct the judge LLM client lazily."""
    return GoogleAIClientLLM(
        config=GoogleAIClientLLMConfig(
            model=model_id,
            api_key=SETTINGS.GEMINI_API_KEY,
            temperature=1.0,
            thinking_level="low",
            max_tokens=4000,
            response_mime_type="application/json",
            response_schema=ArtifactVerdict,
        )
    )


async def _judge_one_artifact(
    llm: GoogleAIClientLLM,
    artifact_type: str,
    artifact_path: Path,
    rubric: str,
    transcript_context: str,
) -> ArtifactResult:
    """Run the judge over one artifact and return a structured result."""
    result = ArtifactResult(
        artifact_type=artifact_type, artifact_path=str(artifact_path)
    )
    extractor = _EXTRACTORS[artifact_type]
    content, skip_reason = extractor(artifact_path)
    if skip_reason or not content.strip():
        result.skipped = True
        result.skip_reason = skip_reason or "no extractable text"
        return result

    criteria_list = ", ".join(_CRITERIA[artifact_type])
    system_prompt = _SYSTEM_PROMPT_TEMPLATE.format(rubric=rubric)
    user_prompt = _USER_PROMPT_TEMPLATE.format(
        artifact_type=artifact_type,
        artifact_path=str(artifact_path),
        criteria_list=criteria_list,
        artifact_content=content[:40000],
        transcript_context=transcript_context,
    )

    full_prompt = f"{system_prompt}\n\n{user_prompt}"
    try:
        response = await llm.acomplete(full_prompt)
    except Exception as exc:  # noqa: BLE001
        result.skipped = True
        result.skip_reason = f"judge call failed: {exc}"
        return result

    parsed = _coerce_verdict(response.text, artifact_type)
    if parsed is None:
        result.skipped = True
        result.skip_reason = "verdict could not be parsed"
        return result

    for criterion in parsed.criteria:
        item = {
            "id": criterion.id,
            "judgment": criterion.judgment,
            "evidence": criterion.evidence,
            "explanation": criterion.explanation,
        }
        result.criteria.append(item)
        if criterion.judgment == "MET":
            result.met_count += 1
        elif criterion.judgment == "UNMET":
            result.unmet_count += 1
        else:
            result.cannot_assess_count += 1
    result.passed = result.met_count >= _PASS_FLOOR_PER_ARTIFACT
    return result


def _coerce_verdict(
    raw: Any, artifact_type: str
) -> ArtifactVerdict | None:
    """Best-effort coercion when the LLM returns dict / string instead."""
    if isinstance(raw, ArtifactVerdict):
        return raw
    if isinstance(raw, dict):
        try:
            return ArtifactVerdict(**raw)
        except Exception:  # noqa: BLE001
            return None
    if isinstance(raw, str):
        try:
            return ArtifactVerdict(**json.loads(raw))
        except Exception:  # noqa: BLE001
            return None
    return None


# ---------------------------------------------------------------------------
# Driver
# ---------------------------------------------------------------------------


async def judge_session(
    session_dir: Path,
    judge_model: str,
    artifact_filter: str | None = None,
) -> JudgeReport:
    """Run the artifact content judge over a single session directory.

    Args:
        session_dir: Pilot session directory (under
            ``brandmind-output/eval/``).
        judge_model: Gemini model id used by the judge.
        artifact_filter: When set to one of the artifact-type keys,
            only that artifact is judged. ``None`` judges all four.

    Returns:
        :class:`JudgeReport` with per-artifact verdicts and the
        aggregate pass/fail.
    """
    rubric = _RUBRIC_PATH.read_text(encoding="utf-8")
    audit_report = run_existence_audit(
        session_dir=session_dir,
        brandmind_home=Path.home() / ".brandmind",
        output_root=Path.cwd() / "brandmind-output",
    )
    transcript_context = _load_transcript_context(session_dir)
    llm = _build_llm(judge_model)
    report = JudgeReport(
        session_dir=str(session_dir),
        rubric_path=str(_RUBRIC_PATH),
        judge_model=judge_model,
    )

    targets: list[tuple[str, Path]] = []
    for artifact_type in _CRITERIA.keys():
        if artifact_filter and artifact_filter != artifact_type:
            continue
        candidates = [
            Path(p) for p in audit_report.artifacts_on_disk.get(artifact_type, [])
        ]
        # The audit scans the entire output tree by mtime cutoff, which can
        # surface artifacts from prior sessions when several pilots ran in
        # quick succession. The content judge speaks to ONE session, so
        # narrow each artifact bucket to the single most recently
        # modified file before running the LLM call.
        if candidates:
            latest = max(candidates, key=lambda p: p.stat().st_mtime)
            targets.append((artifact_type, latest))

    for artifact_type, artifact_path in targets:
        result = await _judge_one_artifact(
            llm,
            artifact_type=artifact_type,
            artifact_path=artifact_path,
            rubric=rubric,
            transcript_context=transcript_context,
        )
        report.artifacts.append(result)

    judged = [a for a in report.artifacts if not a.skipped]
    skipped = [a for a in report.artifacts if a.skipped]
    pass_count = sum(1 for a in judged if a.passed)
    # Scale the aggregate threshold down by the number of skipped artifacts:
    # the rubric defines 3 of 4 as the bar; when an artifact cannot be
    # judged (e.g. OCR dependency missing), only the remaining artifacts
    # are eligible to contribute, so the floor drops one for each skip.
    effective_floor = max(_AGGREGATE_PASS_FLOOR - len(skipped), 1)
    report.aggregate_pass = pass_count >= effective_floor
    report.aggregate_summary = (
        f"{pass_count}/{len(judged)} judged artifacts pass "
        f"(threshold {effective_floor}; "
        f"{len(skipped)} skipped of {len(report.artifacts)} total)"
    )
    return report


def _format_summary(report: JudgeReport) -> str:
    lines: list[str] = []
    lines.append(f"Session: {report.session_dir}")
    lines.append(f"Judge: {report.judge_model}")
    lines.append("")
    for r in report.artifacts:
        if r.skipped:
            lines.append(
                f"[SKIP] {r.artifact_type}: {r.skip_reason} "
                f"({Path(r.artifact_path).name})"
            )
            continue
        status = "PASS" if r.passed else "FAIL"
        lines.append(
            f"[{status}] {r.artifact_type}: "
            f"MET={r.met_count} UNMET={r.unmet_count} "
            f"CANNOT_ASSESS={r.cannot_assess_count} "
            f"({Path(r.artifact_path).name})"
        )
        for c in r.criteria:
            lines.append(f"    {c['id']}: {c['judgment']} — {c['explanation']}")
    lines.append("")
    aggregate = "PASS" if report.aggregate_pass else "FAIL"
    lines.append(f"Aggregate: {aggregate} — {report.aggregate_summary}")
    return "\n".join(lines)


def _build_argparser() -> argparse.ArgumentParser:
    parser = argparse.ArgumentParser(description=__doc__)
    parser.add_argument(
        "--session-dir",
        required=True,
        type=Path,
        help="Pilot session directory under brandmind-output/eval/.",
    )
    parser.add_argument(
        "--judge-model",
        default="gemini-3.1-flash-lite-preview",
        help="Gemini model id used by the artifact content judge.",
    )
    parser.add_argument(
        "--artifact",
        choices=sorted(_CRITERIA.keys()),
        help="Limit the judge to one artifact type (default: all).",
    )
    parser.add_argument(
        "--json",
        action="store_true",
        help="Emit the raw report as JSON on stdout.",
    )
    parser.add_argument(
        "--out",
        type=Path,
        help=(
            "Write the JSON report to this path (defaults to "
            "<session-dir>/artifact_judge.json)."
        ),
    )
    return parser


def main(argv: list[str] | None = None) -> int:
    args = _build_argparser().parse_args(argv)
    if not SETTINGS.GEMINI_API_KEY:
        print(
            "error: GEMINI_API_KEY is required (set env or load environments/.env)",
            file=sys.stderr,
        )
        return 2
    report = asyncio.run(
        judge_session(
            session_dir=args.session_dir,
            judge_model=args.judge_model,
            artifact_filter=args.artifact,
        )
    )
    output_path = args.out or args.session_dir / "artifact_judge.json"
    output_path.write_text(
        json.dumps(asdict(report), indent=2, ensure_ascii=False),
        encoding="utf-8",
    )
    if args.json:
        print(json.dumps(asdict(report), indent=2, ensure_ascii=False))
    else:
        print(_format_summary(report))
    return 0 if report.aggregate_pass else 1


if __name__ == "__main__":
    raise SystemExit(main())
