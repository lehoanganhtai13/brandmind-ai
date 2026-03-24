"""PPTX builder for brand strategy pitch decks using python-pptx.

Generates 10-15 slide pitch decks with 5 layout types:
title, content, two_column, image, table.
"""

from __future__ import annotations

from pathlib import Path
from typing import Any

from loguru import logger
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt

from .templates.brand_strategy import BrandStrategyDeckTemplate


def _hex_to_pptx_rgb(hex_color: str) -> RGBColor:
    """Convert hex color to python-pptx RGBColor."""
    hex_color = hex_color.lstrip("#")
    return RGBColor(
        int(hex_color[0:2], 16),
        int(hex_color[2:4], 16),
        int(hex_color[4:6], 16),
    )


class BrandStrategyPPTXBuilder:
    """Builds professional branded PPTX pitch decks.

    Uses python-pptx with 16:9 widescreen format and
    brand color theming.
    """

    def __init__(self, template: BrandStrategyDeckTemplate | None = None) -> None:
        self.template = template or BrandStrategyDeckTemplate()
        self.primary_color = _hex_to_pptx_rgb(self.template.brand_colors[0])
        self.accent_color = (
            _hex_to_pptx_rgb(self.template.brand_colors[2])
            if len(self.template.brand_colors) > 2
            else _hex_to_pptx_rgb("#D4A84B")
        )

    def build(
        self,
        content: dict[str, Any],
        output_path: str,
        images: dict[str, str] | None = None,
    ) -> str:
        """Build PPTX pitch deck from content dict.

        Args:
            content: Dict with slide content keyed by content_key.
            output_path: File path for output PPTX.
            images: Optional dict mapping slide_id to image path.

        Returns:
            Path to generated PPTX file.
        """
        prs = Presentation()
        prs.slide_width = Inches(13.333)
        prs.slide_height = Inches(7.5)

        for slide_tmpl in self.template.slides:
            slide_content = self._resolve_content(content, slide_tmpl.content_key)

            if images and slide_tmpl.id in images:
                slide_content = images[slide_tmpl.id]

            if slide_content is None and not slide_tmpl.required:
                continue

            layout = slide_tmpl.layout
            if layout == "title":
                self._create_title_slide(prs, slide_tmpl.title, slide_content)
            elif layout == "content":
                self._create_content_slide(
                    prs, slide_tmpl.title, slide_content, slide_tmpl.notes
                )
            elif layout == "two_column":
                self._create_two_column_slide(
                    prs, slide_tmpl.title, slide_content, slide_tmpl.notes
                )
            elif layout == "image":
                self._create_image_slide(
                    prs, slide_tmpl.title, slide_content, slide_tmpl.notes
                )
            elif layout == "table":
                self._create_table_slide(
                    prs, slide_tmpl.title, slide_content, slide_tmpl.notes
                )

        Path(output_path).parent.mkdir(parents=True, exist_ok=True)
        prs.save(output_path)
        logger.info(f"PPTX generated: {output_path}")
        return output_path

    def _create_title_slide(self, prs: Any, title: str, content: Any) -> None:
        """Create title/cover slide."""
        layout = prs.slide_layouts[0]
        slide = prs.slides.add_slide(layout)

        slide.shapes.title.text = self.template.brand_name
        for run in slide.shapes.title.text_frame.paragraphs[0].runs:
            run.font.color.rgb = self.primary_color
            run.font.size = Pt(44)

        if slide.placeholders.get(1):
            slide.placeholders[1].text = title

    def _create_content_slide(
        self,
        prs: Any,
        title: str,
        content: Any,
        notes: str = "",
    ) -> None:
        """Create title + bullet points slide."""
        layout = prs.slide_layouts[1]
        slide = prs.slides.add_slide(layout)

        slide.shapes.title.text = title
        for run in slide.shapes.title.text_frame.paragraphs[0].runs:
            run.font.color.rgb = self.primary_color

        bullets = self._extract_bullets(content)
        body = slide.placeholders.get(1)
        if body:
            tf = body.text_frame
            tf.clear()
            for i, bullet in enumerate(bullets[:10]):
                if i == 0:
                    tf.paragraphs[0].text = bullet
                else:
                    p = tf.add_paragraph()
                    p.text = bullet
                    p.font.size = Pt(14)

        if notes:
            slide.notes_slide.notes_text_frame.text = notes

    def _create_two_column_slide(
        self,
        prs: Any,
        title: str,
        content: Any,
        notes: str = "",
    ) -> None:
        """Create side-by-side comparison layout."""
        layout = prs.slide_layouts[1]
        slide = prs.slides.add_slide(layout)

        slide.shapes.title.text = title
        for run in slide.shapes.title.text_frame.paragraphs[0].runs:
            run.font.color.rgb = self.primary_color

        bullets = self._extract_bullets(content)
        body = slide.placeholders.get(1)
        if body:
            tf = body.text_frame
            tf.clear()
            for i, bullet in enumerate(bullets[:10]):
                if i == 0:
                    tf.paragraphs[0].text = bullet
                else:
                    p = tf.add_paragraph()
                    p.text = bullet

        if notes:
            slide.notes_slide.notes_text_frame.text = notes

    def _create_image_slide(
        self,
        prs: Any,
        title: str,
        content: Any,
        notes: str = "",
    ) -> None:
        """Create image slide with title."""
        layout = prs.slide_layouts[6]  # Blank
        slide = prs.slides.add_slide(layout)

        txBox = slide.shapes.add_textbox(
            Inches(0.5), Inches(0.3), Inches(12), Inches(0.8)
        )
        tf = txBox.text_frame
        p = tf.paragraphs[0]
        p.text = title
        p.font.size = Pt(28)
        p.font.bold = True
        p.font.color.rgb = self.primary_color

        image_path = content if isinstance(content, str) else None
        if image_path and Path(image_path).exists():
            slide.shapes.add_picture(
                image_path,
                Inches(1.5),
                Inches(1.3),
                Inches(10),
                Inches(5.5),
            )
        else:
            txBox2 = slide.shapes.add_textbox(
                Inches(2), Inches(3), Inches(9), Inches(2)
            )
            tf2 = txBox2.text_frame
            tf2.paragraphs[0].text = "(Image placeholder)"
            tf2.paragraphs[0].alignment = PP_ALIGN.CENTER

        if notes:
            slide.notes_slide.notes_text_frame.text = notes

    def _create_table_slide(
        self,
        prs: Any,
        title: str,
        content: Any,
        notes: str = "",
    ) -> None:
        """Create data table slide."""
        layout = prs.slide_layouts[6]  # Blank
        slide = prs.slides.add_slide(layout)

        txBox = slide.shapes.add_textbox(
            Inches(0.5), Inches(0.3), Inches(12), Inches(0.8)
        )
        tf = txBox.text_frame
        p = tf.paragraphs[0]
        p.text = title
        p.font.size = Pt(28)
        p.font.bold = True
        p.font.color.rgb = self.primary_color

        if isinstance(content, list) and content and isinstance(content[0], dict):
            headers = list(content[0].keys())[:6]
            rows = content[:8]
            n_rows = len(rows) + 1
            n_cols = len(headers)

            table_shape = slide.shapes.add_table(
                n_rows,
                n_cols,
                Inches(0.5),
                Inches(1.3),
                Inches(12),
                Inches(5.5),
            )
            table = table_shape.table

            for i, header in enumerate(headers):
                cell = table.cell(0, i)
                cell.text = header.replace("_", " ").title()
                for paragraph in cell.text_frame.paragraphs:
                    for run in paragraph.runs:
                        run.font.bold = True
                        run.font.color.rgb = RGBColor(255, 255, 255)
                        run.font.size = Pt(10)
                cell.fill.solid()
                cell.fill.fore_color.rgb = self.primary_color

            for row_idx, row_data in enumerate(rows, 1):
                for col_idx, header in enumerate(headers):
                    cell = table.cell(row_idx, col_idx)
                    cell.text = str(row_data.get(header, ""))
        else:
            bullets = self._extract_bullets(content)
            txBox2 = slide.shapes.add_textbox(
                Inches(0.5), Inches(1.3), Inches(12), Inches(5.5)
            )
            tf2 = txBox2.text_frame
            for i, bullet in enumerate(bullets[:10]):
                if i == 0:
                    tf2.paragraphs[0].text = f"• {bullet}"
                else:
                    p = tf2.add_paragraph()
                    p.text = f"• {bullet}"

        if notes:
            slide.notes_slide.notes_text_frame.text = notes

    def _extract_bullets(self, content: Any) -> list[str]:
        """Convert content to bullet-point list."""
        if isinstance(content, str):
            lines = [line.strip() for line in content.split("\n") if line.strip()]
            return lines[:10]
        elif isinstance(content, list):
            return [str(item) for item in content[:10]]
        elif isinstance(content, dict):
            bullets = []
            for key, value in content.items():
                label = key.replace("_", " ").title()
                if isinstance(value, str):
                    bullets.append(f"{label}: {value}")
                elif isinstance(value, list):
                    bullets.append(f"{label}: {', '.join(str(v) for v in value[:3])}")
                else:
                    bullets.append(f"{label}: {str(value)[:80]}")
            return bullets[:10]
        return ["(No content)"]

    def _resolve_content(self, content: dict[str, Any], content_key: str) -> Any:
        """Resolve dotted content key path."""
        parts = content_key.split(".")
        current = content
        for part in parts:
            if isinstance(current, dict) and part in current:
                current = current[part]
            else:
                return None
        return current
